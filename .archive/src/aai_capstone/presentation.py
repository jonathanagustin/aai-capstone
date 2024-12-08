# src/aai_capstone/presentation.py
from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.dml.color import RGBColor
from . import OUTPUT_DIR
from .config import TemplateConfig


class PresentationBuilder:
    """Build PowerPoint presentations based on a template and YAML content."""

    def __init__(self, template_path: Path, config: TemplateConfig):
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        self.prs: PresentationType = Presentation(template_path)
        self.config = config
        self.usd_blue: RGBColor = RGBColor(0, 59, 112)
        self._create_layout_map()

    def _create_layout_map(self) -> None:
        """Create a map of layout names to their indices for quick access."""
        self.layout_map = {layout.name: idx for idx, layout in enumerate(self.prs.slide_layouts)}

    def get_layout_index(self, name: str) -> int:
        """Get the index of a layout by name.

        Args:
            name (str): Layout name.

        Returns:
            int: Index of the layout.
        """
        if name not in self.layout_map:
            raise ValueError(f"Layout '{name}' not found in template")
        return self.layout_map[name]

    def create_section_header(self, title: str) -> None:
        """Create a section header slide with a given title."""
        layout_idx = self.get_layout_index("1 - Section Header")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

        layout_config = self.config.get_layout_by_name("1 - Section Header")
        title_placeholder = next(p for p in layout_config['placeholders'] if p['type_id'] == 'TITLE')
        title_shape = slide.placeholders[title_placeholder['index']]
        title_shape.text = title

    def create_content_slide(self, title: str, content: List[Dict[str, Any]]) -> None:
        """Create a content slide with bullet points."""
        layout_idx = self.get_layout_index("2 - Title and Content")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

        layout_config = self.config.get_layout_by_name("2 - Title and Content")

        # Set title
        title_placeholder = next(p for p in layout_config['placeholders'] if p['type_id'] == 'TITLE')
        title_shape = slide.placeholders[title_placeholder['index']]
        title_shape.text = title

        # Set content
        body_placeholder = next(p for p in layout_config['placeholders'] if p['type_id'] == 'BODY')
        content_shape = slide.placeholders[body_placeholder['index']]
        self._add_content_to_shape(content_shape, content)

    def create_comparison_slide(
        self,
        title: str,
        left_content: List[Dict[str, Any]],
        right_content: List[Dict[str, Any]]
    ) -> None:
        """Create a comparison slide with two columns of bullet points."""
        layout_idx = self.get_layout_index("3 - Comparison")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

        layout_config = self.config.get_layout_by_name("3 - Comparison")

        # Set title
        title_placeholder = next(p for p in layout_config['placeholders'] if p['type_id'] == 'TITLE')
        title_shape = slide.placeholders[title_placeholder['index']]
        title_shape.text = title

        # Set left content
        left_body = next(p for p in layout_config['placeholders'] if p['type_id'] == 'BODY' and p['index'] == 1)
        left_shape = slide.placeholders[left_body['index']]
        self._add_content_to_shape(left_shape, left_content)

        # Set right content
        right_body = next(p for p in layout_config['placeholders'] if p['type_id'] == 'BODY' and p['index'] == 13)
        right_shape = slide.placeholders[right_body['index']]
        self._add_content_to_shape(right_shape, right_content)

    def _add_content_to_shape(self, shape, content: List[Dict[str, Any]]) -> None:
        """Add bullet points to a shape's text frame.

        Args:
            shape: The placeholder shape to which text will be added.
            content (List[Dict[str, Any]]): A list of items or dicts with 'text' and optional 'level'.
        """
        tf = shape.text_frame
        tf.clear()

        for item in content:
            p = tf.add_paragraph()
            if isinstance(item, dict):
                p.text = item['text']
                p.level = item.get('level', 0)
            else:
                p.text = str(item)
                p.level = 0

    def save(self, filename: Path) -> None:
        """Save the built presentation to a PPTX file."""
        self.prs.save(filename)
