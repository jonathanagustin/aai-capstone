#!/usr/bin/env python3
# src/aai_capstone/analyze.py
from __future__ import annotations

import os
import sys
from typing import Dict, Any, NoReturn
import yaml

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.shapes.picture import CT_Picture
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.oxml.shapes.autoshape import CT_Shape
from pptx.oxml.ns import qn

class TemplateAnalyzer:
    """Analyze a PowerPoint template to determine layouts, placeholders, and shapes with more detail."""

    def ensure_config_directory(self) -> None:
        os.makedirs("config", exist_ok=True)

    def analyze_template(self, template_path: str) -> Dict[str, Any]:
        if not os.path.isfile(template_path):
            print(f"Error: Template not found at {template_path}", file=sys.stderr)
            sys.exit(1)

        try:
            prs = Presentation(template_path)
        except Exception as e:
            print(f"Error loading template: {e}", file=sys.stderr)
            sys.exit(1)

        template_config = {
            "template": {
                "path": template_path,
                "dimensions": {
                    "width": round(prs.slide_width / 914400, 2),
                    "height": round(prs.slide_height / 914400, 2)
                }
            },
            "layouts": [],
        }

        total_layouts = len(prs.slide_layouts)
        total_placeholders = 0
        total_shapes = 0

        for layout_idx, layout in enumerate(prs.slide_layouts, start=1):
            layout_name = layout.name or f"Layout {layout_idx}"
            shape_records = []
            placeholder_count = 0
            # Gather shapes
            for shape in layout.shapes:
                shape_info = self._shape_info(shape)
                shape_records.append(shape_info)
                total_shapes += 1
                if shape_info.get('placeholder') is not None:
                    placeholder_count += 1
                    total_placeholders += 1

            layout_config = {
                "name": layout_name,
                "layout_index": layout_idx,
                "placeholder_count": placeholder_count,
                "shape_count": len(shape_records),
                "shapes": shape_records
            }
            template_config["layouts"].append(layout_config)

        template_config["summary"] = {
            "total_layouts": total_layouts,
            "total_placeholders": total_placeholders,
            "total_shapes": total_shapes
        }

        return template_config

    def _shape_info(self, shape) -> Dict[str, Any]:
        """Return a dictionary of attributes about a shape, including placeholder info, shape type, and text."""
        shape_info = {}
        shape_info["shape_id"] = getattr(shape, 'shape_id', None)
        shape_info["name"] = shape.name if hasattr(shape, 'name') else "<No Name>"

        # Determine shape_type (placeholder, picture, graphicframe, etc.)
        elm = shape.element
        if elm.tag == qn("p:sp"):
            # autoshape or placeholder
            if shape.is_placeholder:
                shape_info["shape_type"] = "placeholder"
                placeholder_info = self._placeholder_info(shape)
                shape_info["placeholder"] = placeholder_info
            else:
                shape_info["shape_type"] = "autoshape"
        elif elm.tag == qn("p:pic"):
            shape_info["shape_type"] = "picture"
            # Could extract image desc or rId from elm if needed
            shape_info["image"] = {"desc": elm.xpath('.//p:nvPicPr/p:cNvPr/@descr') or ""}
        elif elm.tag == qn("p:graphicFrame"):
            shape_info["shape_type"] = "graphicFrame"
            # Check if this is a chart or a table
            if elm.xpath('.//c:chart', namespaces=elm.nsmap):
                shape_info["graphic_type"] = "chart"
            elif elm.xpath('.//a:tbl', namespaces=elm.nsmap):
                shape_info["graphic_type"] = "table"
            else:
                shape_info["graphic_type"] = "unknown_graphic"
        elif elm.tag == qn("p:cxnSp"):
            shape_info["shape_type"] = "connector"
        elif elm.tag == qn("p:grpSp"):
            shape_info["shape_type"] = "group"
            # could count shapes inside group if desired
        else:
            shape_info["shape_type"] = "unknown"

        # Extract text if applicable
        if hasattr(shape, 'text') and shape.text:
            shape_info["text"] = shape.text.strip()

        return shape_info

    def _placeholder_info(self, shape) -> Dict[str, Any]:
        """Return dictionary with placeholder details if shape is a placeholder."""
        ph = shape.element.ph
        if ph is None:
            return {}
        # ph_type, sz, orient, idx
        ph_info = {
            "ph_type": str(ph.type) if ph.type else None,
            "ph_idx": ph.idx,
            "ph_sz": ph.sz,
            "ph_orient": ph.orient
        }
        return ph_info

    def save_config(self, config: Dict[str, Any], output_path: str = "config/template.yml") -> None:
        self.ensure_config_directory()
        try:
            with open(output_path, "w", encoding="utf-8") as f:
                yaml.dump(config, f, default_flow_style=False, allow_unicode=True, sort_keys=False)
            print(f"\nConfiguration saved to {output_path}")
        except Exception as e:
            print(f"Error saving configuration: {e}", file=sys.stderr)
            sys.exit(1)

    def load_config(self, config_path: str = "config/template.yml") -> Dict[str, Any]:
        if not os.path.exists(config_path):
            print(f"Error: Configuration file not found at {config_path}", file=sys.stderr)
            sys.exit(1)
        try:
            with open(config_path, "r", encoding="utf-8") as f:
                return yaml.safe_load(f)
        except Exception as e:
            print(f"Error loading configuration: {e}", file=sys.stderr)
            sys.exit(1)

    def print_config(self, config: Dict[str, Any]) -> None:
        print("\nTemplate Analysis")
        print("=================")
        print("\nSlide Dimensions:")
        print(f"  Width:  {config['template']['dimensions']['width']:.2f} inches")
        print(f"  Height: {config['template']['dimensions']['height']:.2f} inches")

        summary = config.get("summary", {})
        print("\nSummary:")
        print("--------")
        print(f"  Total Layouts: {summary.get('total_layouts', 'N/A')}")
        print(f"  Total Placeholders: {summary.get('total_placeholders', 'N/A')}")
        print(f"  Total Shapes: {summary.get('total_shapes', 'N/A')}")

        layouts = config.get("layouts", [])
        print("\nAvailable Layouts:")
        print("-----------------")
        if not layouts:
            print("  No layouts found.")
            return

        for layout_data in layouts:
            print(f"\nLayout {layout_data['layout_index']}: \"{layout_data['name']}\"")
            print(f"  Placeholder Count: {layout_data['placeholder_count']}")
            print(f"  Shape Count: {layout_data['shape_count']}")
            for shape in layout_data["shapes"]:
                print("    -")
                print(f"      shape_id:   {shape.get('shape_id')}")
                print(f"      name:       {shape.get('name')}")
                print(f"      shape_type: {shape.get('shape_type')}")
                if "placeholder" in shape and shape["placeholder"]:
                    print("      placeholder:")
                    for k, v in shape["placeholder"].items():
                        print(f"        {k}: {v}")
                if "image" in shape:
                    print("      image:")
                    for k, v in shape["image"].items():
                        print(f"        {k}: {v}")
                if "text" in shape:
                    print(f"      text: \"{shape['text']}\"")

def main() -> NoReturn:
    analyzer = TemplateAnalyzer()

    if len(sys.argv) < 2:
        print("Usage:")
        print("  python analyze_template.py analyze <template_path>")
        print("  python analyze_template.py load")
        sys.exit(1)

    command = sys.argv[1]

    if command == "analyze":
        if len(sys.argv) < 3:
            print("Error: Template path required for analysis", file=sys.stderr)
            sys.exit(1)

        template_path = sys.argv[2]
        config = analyzer.analyze_template(template_path)
        analyzer.save_config(config)
        analyzer.print_config(config)

    elif command == "load":
        config = analyzer.load_config()
        analyzer.print_config(config)
    else:
        print(f"Unknown command: {command}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
