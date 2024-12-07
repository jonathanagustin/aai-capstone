#!/usr/bin/env python3
from __future__ import annotations

import os
import sys
import yaml
from typing import List, Dict, Any, Union
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.dml.color import RGBColor

class SlideCreator:
    def __init__(self, template_path: str = "assets/templates/template.pptx") -> None:
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")

        self.prs: PresentationType = Presentation(template_path)
        self.usd_blue: RGBColor = RGBColor(0, 59, 112)
        self.config = self.load_template_config()
        self.layout_map = self._create_layout_map()

    def load_template_config(self) -> Dict[str, Any]:
        config_path = "config/template.yml"
        if not os.path.exists(config_path):
            raise FileNotFoundError(f"Template configuration not found: {config_path}")
        with open(config_path, 'r', encoding='utf-8') as f:
            return yaml.safe_load(f)

    def _create_layout_map(self) -> Dict[str, int]:
        return {layout.name: idx for idx, layout in enumerate(self.prs.slide_layouts)}

    def get_layout_by_name(self, name: str) -> int:
        if name not in self.layout_map:
            raise ValueError(f"Layout '{name}' not found in template")
        return self.layout_map[name]

    def process_content_item(self, item: Union[str, Dict[str, Any]]) -> tuple[str, int]:
        if isinstance(item, str):
            return item, 0
        elif isinstance(item, dict):
            return item['text'], item.get('level', 0)
        else:
            raise ValueError(f"Invalid content item format: {item}")

    def create_section_header(self, title: str) -> None:
        layout_idx = self.get_layout_by_name("1 - Section Header")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

        layout_config = next(l for l in self.config['layouts']
                           if l['name'] == "1 - Section Header")
        title_placeholder = next(p for p in layout_config['placeholders']
                               if p['type_id'] == 'TITLE')
        title_shape = slide.placeholders[title_placeholder['index']]
        title_shape.text = title

    def create_content_slide(self, title: str, content: List[Union[str, Dict[str, Any]]]) -> None:
        layout_idx = self.get_layout_by_name("2 - Title and Content")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

        layout_config = next(l for l in self.config['layouts']
                           if l['name'] == "2 - Title and Content")

        # Set title
        title_placeholder = next(p for p in layout_config['placeholders']
                               if p['type_id'] == 'TITLE')
        title_shape = slide.placeholders[title_placeholder['index']]
        title_shape.text = title

        # Set content
        body_placeholder = next(p for p in layout_config['placeholders']
                              if p['type_id'] == 'BODY')
        content_shape = slide.placeholders[body_placeholder['index']]
        tf = content_shape.text_frame
        tf.clear()

        for item in content:
            text, level = self.process_content_item(item)
            p = tf.add_paragraph()
            p.text = text
            p.level = level

    def create_comparison_slide(self, title: str,
                              left_content: List[Union[str, Dict[str, Any]]],
                              right_content: List[Union[str, Dict[str, Any]]]) -> None:
        layout_idx = self.get_layout_by_name("3 - Comparison")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])

        layout_config = next(l for l in self.config['layouts']
                           if l['name'] == "3 - Comparison")

        # Set title
        title_placeholder = next(p for p in layout_config['placeholders']
                               if p['type_id'] == 'TITLE')
        title_shape = slide.placeholders[title_placeholder['index']]
        title_shape.text = title

        # Set left content
        left_body = next(p for p in layout_config['placeholders']
                        if p['type_id'] == 'BODY' and p['index'] == 1)
        left_shape = slide.placeholders[left_body['index']]
        self._add_content_to_shape(left_shape, left_content)

        # Set right content
        right_body = next(p for p in layout_config['placeholders']
                         if p['type_id'] == 'BODY' and p['index'] == 13)
        right_shape = slide.placeholders[right_body['index']]
        self._add_content_to_shape(right_shape, right_content)

    def _add_content_to_shape(self, shape, content: List[Union[str, Dict[str, Any]]]) -> None:
        tf = shape.text_frame
        tf.clear()

        for item in content:
            text, level = self.process_content_item(item)
            p = tf.add_paragraph()
            p.text = text
            p.level = level

    def create_presentation_from_yaml(self, content_path: str) -> None:
        with open(content_path, 'r', encoding='utf-8') as f:
            content = yaml.safe_load(f)

        for slide in content['slides']:
            slide_type = slide['type']

            if slide_type == 'section':
                self.create_section_header(slide['title'])
            elif slide_type == 'content':
                self.create_content_slide(slide['title'], slide['content'])
            elif slide_type == 'comparison':
                self.create_comparison_slide(
                    slide['title'],
                    slide['left_content'],
                    slide['right_content']
                )
            else:
                raise ValueError(f"Unknown slide type: {slide_type}")

    def save(self, filename: str) -> None:
        directory = os.path.dirname(filename)
        if directory:
            os.makedirs(directory, exist_ok=True)
        self.prs.save(filename)

def main() -> None:
    try:
        creator = SlideCreator()

        # Process all YAML files in the content directory
        content_dir = "content"
        if not os.path.exists(content_dir):
            raise FileNotFoundError(f"Content directory not found: {content_dir}")

        for filename in sorted(os.listdir(content_dir)):
            if filename.endswith('.yml'):
                print(f"Processing {filename}...")
                content_path = os.path.join(content_dir, filename)
                output_name = f"{os.path.splitext(filename)[0]}.pptx"

                creator = SlideCreator()  # Fresh instance for each presentation
                creator.create_presentation_from_yaml(content_path)
                creator.save(output_name)
                print(f"Created presentation: {output_name}")

    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()