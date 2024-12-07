#!/usr/bin/env python3
# scripts/analyze_template.py

from __future__ import annotations

import os
import sys
from typing import Dict, Any, NoReturn
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.enum.shapes import PP_PLACEHOLDER
import yaml


class TemplateAnalyzer:
    def ensure_config_directory(self) -> None:
        """Ensure the config directory exists."""
        os.makedirs("config", exist_ok=True)

    def get_placeholder_type(self, ph_type: int) -> tuple[str, str]:
        """Convert placeholder type to string representation."""
        try:
            ph_enum = PP_PLACEHOLDER(ph_type)
            type_str = f"{ph_enum.name} ({ph_type})"
            type_id = f"{ph_enum.name}"
            return type_str, type_id
        except ValueError:
            return f"Type {ph_type}", str(ph_type)

    def analyze_template(self, template_path: str) -> Dict[str, Any]:
        """Analyze PowerPoint template and return configuration dictionary."""
        if not os.path.exists(template_path):
            print(f"Error: Template not found at {template_path}")
            sys.exit(1)

        try:
            prs: PresentationType = Presentation(template_path)
        except Exception as e:
            print(f"Error loading template: {e}")
            sys.exit(1)

        config = {
            "template": {
                "path": template_path,
                "dimensions": {
                    "width": round(prs.slide_width / 914400, 2),  # Convert EMU to inches
                    "height": round(prs.slide_height / 914400, 2)
                }
            },
            "layouts": []
        }

        # Analyze layouts
        for layout in prs.slide_layouts:
            layout_config = {
                "name": layout.name or "",  # Handle blank name
                "placeholders": []
            }

            for shape in layout.placeholders:
                ph_type = shape.placeholder_format.type
                type_str, type_id = self.get_placeholder_type(ph_type)

                placeholder_config = {
                    "index": shape.placeholder_format.idx if hasattr(shape.placeholder_format, 'idx') else None,
                    "type": type_str or "UNKNOWN",  # Use "UNKNOWN" for blank types
                    "type_id": type_id or "",  # Empty string for blank type_id
                    "name": shape.name or ""  # Handle blank name
                }
                layout_config["placeholders"].append(placeholder_config)

            config["layouts"].append(layout_config)

        return config

    def save_config(self, config: Dict[str, Any], output_path: str = "config/template.yml") -> None:
        """Save configuration to YAML file using PyYAML."""
        self.ensure_config_directory()

        try:
            with open(output_path, "w", encoding="utf-8") as f:
                yaml.dump(config, f, default_flow_style=False, allow_unicode=True, sort_keys=False)
            print(f"\nConfiguration saved to {output_path}")
        except Exception as e:
            print(f"Error saving configuration: {e}")
            sys.exit(1)

    def load_config(self, config_path: str = "config/template.yml") -> Dict[str, Any]:
        """Load configuration from YAML file."""
        if not os.path.exists(config_path):
            print(f"Error: Configuration file not found at {config_path}")
            sys.exit(1)
        try:
            with open(config_path, "r", encoding="utf-8") as f:
                return yaml.safe_load(f)
        except Exception as e:
            print(f"Error loading configuration: {e}")
            sys.exit(1)

    def print_config(self, config: Dict[str, Any]) -> None:
        """Print configuration in a readable format."""
        print("\nTemplate Analysis")
        print("=================")

        print("\nSlide Dimensions:")
        print(f"Width:  {config['template']['dimensions']['width']:.2f} inches")
        print(f"Height: {config['template']['dimensions']['height']:.2f} inches")

        print("\nAvailable Layouts:")
        print("-----------------")

        for layout_data in config["layouts"]:
            # Format layout name to add space between number and dash
            layout_name = layout_data['name']
            layout_name = ' '.join(part.strip() for part in layout_name.replace('-', ' -').split())
            print(f"\nLayout: \"{layout_name}\"")
            print("Placeholders:")
            for ph in layout_data["placeholders"]:
                print(f"  - Index: {ph['index']}")
                print(f"    Type:  {ph['type']}")
                print(f"    Name:  {ph['name']}")


def main() -> NoReturn:
    analyzer = TemplateAnalyzer()

    if len(sys.argv) < 2:
        print("Usage:")
        print("  Analyze and save: python analyze_template.py analyze <template_path>")
        print("  Load and display: python analyze_template.py load")
        sys.exit(1)

    command = sys.argv[1]

    if command == "analyze":
        if len(sys.argv) < 3:
            print("Error: Template path required for analysis")
            sys.exit(1)

        template_path = sys.argv[2]
        config = analyzer.analyze_template(template_path)
        analyzer.save_config(config)
        analyzer.print_config(config)

    elif command == "load":
        config = analyzer.load_config()
        analyzer.print_config(config)

    else:
        print(f"Unknown command: {command}")
        sys.exit(1)


if __name__ == "__main__":
    main()
